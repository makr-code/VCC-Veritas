"""
Excel Table Agent - Generate and embed tables in Excel, Word, and PowerPoint

This agent creates formatted tables from data and can embed them in various
document formats following the template-based approach.

Capabilities:
- Generate Excel spreadsheets (.xlsx)
- Generate CSV files
- Embed tables in Word documents
- Embed tables in PowerPoint presentations
- Apply formatting, formulas, and styles
- Support for various table types (data, comparison, summary, schedule)
"""

import logging
import pandas as pd
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime

# Excel/Table libraries
try:
    import openpyxl
    from openpyxl.styles import Font, Fill, Alignment, Border, Side, PatternFill
    from openpyxl.utils.dataframe import dataframe_to_rows
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logging.warning("openpyxl not available - Excel generation disabled")

# Word integration
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available - Word embedding disabled")

# PowerPoint integration
try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.dml.color import RGBColor as PptRGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PowerPoint embedding disabled")

from backend.agents.table_template_manager import get_table_template_manager

logger = logging.getLogger(__name__)


class ExcelTableAgent:
    """
    Agent for Excel/Table Generation and Embedding
    
    Generates formatted tables and embeds them in various document formats.
    """
    
    def __init__(self, llm_service=None, output_dir: Optional[Path] = None):
        """
        Initialize Excel Table Agent
        
        Args:
            llm_service: Optional LLM service for data extraction
            output_dir: Directory for output files
        """
        self.llm_service = llm_service
        self.template_manager = get_table_template_manager()
        
        if output_dir is None:
            self.output_dir = Path("./output/tables")
        else:
            self.output_dir = Path(output_dir)
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info("ExcelTableAgent initialized")
    
    async def generate_table(
        self,
        request: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Generate a table based on request
        
        Args:
            request: {
                'template': template name (e.g., 'data_table'),
                'variation': variation ID (e.g., 'simple_data_table'),
                'data': {
                    'headers': [...],
                    'rows': [[...], [...]],
                    ...
                },
                'output_format': 'excel' | 'csv' | 'word' | 'powerpoint',
                'filename': optional filename
            }
            
        Returns:
            Dictionary with file path and metadata
        """
        try:
            # Step 1: Load template
            template_name = request.get('template', 'data_table')
            variation_id = request.get('variation')
            
            template = self.template_manager.read_template(template_name)
            if not template:
                return {
                    'success': False,
                    'error': f"Template {template_name} not found"
                }
            
            variation = None
            if variation_id:
                variation = self.template_manager.get_variation(template_name, variation_id)
            
            # Step 2: Extract data
            data = request.get('data', {})
            headers = data.get('headers', [])
            rows = data.get('rows', [])
            
            if not headers or not rows:
                return {
                    'success': False,
                    'error': "Missing headers or rows in data"
                }
            
            # Step 3: Create DataFrame
            df = pd.DataFrame(rows, columns=headers)
            
            # Step 4: Generate output based on format
            output_format = request.get('output_format', 'excel')
            filename = request.get('filename', f"table_{datetime.now().strftime('%Y%m%d_%H%M%S')}")
            
            result = None
            if output_format == 'excel':
                result = self._create_excel(df, template, variation, filename)
            elif output_format == 'csv':
                result = self._create_csv(df, filename)
            elif output_format == 'word':
                result = self._create_word_table(df, template, variation, filename)
            elif output_format == 'powerpoint':
                result = self._create_powerpoint_table(df, template, variation, filename)
            else:
                return {
                    'success': False,
                    'error': f"Unsupported output format: {output_format}"
                }
            
            return result
        
        except Exception as e:
            logger.error(f"Table generation failed: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
    
    def _create_excel(
        self,
        df: pd.DataFrame,
        template: Dict[str, Any],
        variation: Optional[Dict[str, Any]],
        filename: str
    ) -> Dict[str, Any]:
        """Create Excel file with formatting"""
        if not OPENPYXL_AVAILABLE:
            return {
                'success': False,
                'error': "openpyxl not available"
            }
        
        try:
            # Create workbook
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Data"
            
            # Write headers
            for col_idx, header in enumerate(df.columns, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                cell.alignment = Alignment(horizontal="center", vertical="center")
            
            # Write data
            for row_idx, row_data in enumerate(df.values, 2):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    
                    # Alternate row coloring
                    if row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Save file
            filepath = self.output_dir / f"{filename}.xlsx"
            wb.save(filepath)
            
            logger.info(f"Excel file created: {filepath}")
            
            return {
                'success': True,
                'file_path': str(filepath),
                'format': 'excel',
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"Excel creation failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _create_csv(self, df: pd.DataFrame, filename: str) -> Dict[str, Any]:
        """Create CSV file"""
        try:
            filepath = self.output_dir / f"{filename}.csv"
            df.to_csv(filepath, index=False, encoding='utf-8-sig')
            
            logger.info(f"CSV file created: {filepath}")
            
            return {
                'success': True,
                'file_path': str(filepath),
                'format': 'csv',
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"CSV creation failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _create_word_table(
        self,
        df: pd.DataFrame,
        template: Dict[str, Any],
        variation: Optional[Dict[str, Any]],
        filename: str
    ) -> Dict[str, Any]:
        """Create Word document with table"""
        if not DOCX_AVAILABLE:
            return {
                'success': False,
                'error': "python-docx not available"
            }
        
        try:
            doc = Document()
            
            # Add title
            title = doc.add_heading('Data Table', level=1)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Create table
            table = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
            table.style = 'Light Grid Accent 1'
            
            # Add headers
            header_cells = table.rows[0].cells
            for idx, header in enumerate(df.columns):
                header_cells[idx].text = str(header)
                # Bold headers
                for paragraph in header_cells[idx].paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            
            # Add data
            for row_idx, row_data in enumerate(df.values, 1):
                row_cells = table.rows[row_idx].cells
                for col_idx, value in enumerate(row_data):
                    row_cells[col_idx].text = str(value)
            
            # Save file
            filepath = self.output_dir / f"{filename}.docx"
            doc.save(filepath)
            
            logger.info(f"Word file created: {filepath}")
            
            return {
                'success': True,
                'file_path': str(filepath),
                'format': 'word',
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"Word creation failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _create_powerpoint_table(
        self,
        df: pd.DataFrame,
        template: Dict[str, Any],
        variation: Optional[Dict[str, Any]],
        filename: str
    ) -> Dict[str, Any]:
        """Create PowerPoint presentation with table"""
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'error': "python-pptx not available"
            }
        
        try:
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = "Data Table"
            
            # Add table slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate table dimensions
            rows = len(df) + 1  # +1 for headers
            cols = len(df.columns)
            
            left = PptInches(1)
            top = PptInches(1.5)
            width = PptInches(8)
            height = PptInches(4.5)
            
            # Add table
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set headers
            for col_idx, header in enumerate(df.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptRGBColor(68, 114, 196)
                
                # Make text white and bold
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = PptRGBColor(255, 255, 255)
            
            # Set data
            for row_idx, row_data in enumerate(df.values, 1):
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value)
                    
                    # Alternate row coloring
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PptRGBColor(242, 242, 242)
            
            # Save file
            filepath = self.output_dir / f"{filename}.pptx"
            prs.save(filepath)
            
            logger.info(f"PowerPoint file created: {filepath}")
            
            return {
                'success': True,
                'file_path': str(filepath),
                'format': 'powerpoint',
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"PowerPoint creation failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def embed_in_existing_word(
        self,
        df: pd.DataFrame,
        document_path: str,
        position: str = "end"
    ) -> Dict[str, Any]:
        """Embed table in existing Word document"""
        if not DOCX_AVAILABLE:
            return {
                'success': False,
                'error': "python-docx not available"
            }
        
        try:
            doc = Document(document_path)
            
            # Create table
            table = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
            table.style = 'Light Grid Accent 1'
            
            # Add headers
            header_cells = table.rows[0].cells
            for idx, header in enumerate(df.columns):
                header_cells[idx].text = str(header)
            
            # Add data
            for row_idx, row_data in enumerate(df.values, 1):
                row_cells = table.rows[row_idx].cells
                for col_idx, value in enumerate(row_data):
                    row_cells[col_idx].text = str(value)
            
            # Save
            doc.save(document_path)
            
            logger.info(f"Table embedded in Word document: {document_path}")
            
            return {
                'success': True,
                'document_path': document_path,
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"Word embedding failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def embed_in_existing_powerpoint(
        self,
        df: pd.DataFrame,
        presentation_path: str,
        slide_index: Optional[int] = None
    ) -> Dict[str, Any]:
        """Embed table in existing PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'error': "python-pptx not available"
            }
        
        try:
            prs = Presentation(presentation_path)
            
            # Add new slide or use existing
            if slide_index is None:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
            else:
                slide = prs.slides[slide_index]
            
            # Add table
            rows = len(df) + 1
            cols = len(df.columns)
            
            left = PptInches(1)
            top = PptInches(1.5)
            width = PptInches(8)
            height = PptInches(4.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set headers and data (similar to _create_powerpoint_table)
            for col_idx, header in enumerate(df.columns):
                table.cell(0, col_idx).text = str(header)
            
            for row_idx, row_data in enumerate(df.values, 1):
                for col_idx, value in enumerate(row_data):
                    table.cell(row_idx, col_idx).text = str(value)
            
            # Save
            prs.save(presentation_path)
            
            logger.info(f"Table embedded in PowerPoint: {presentation_path}")
            
            return {
                'success': True,
                'presentation_path': presentation_path,
                'rows': len(df),
                'columns': len(df.columns)
            }
        
        except Exception as e:
            logger.error(f"PowerPoint embedding failed: {e}")
            return {
                'success': False,
                'error': str(e)
            }


# Standalone test
if __name__ == '__main__':
    import asyncio
    
    async def test_agent():
        agent = ExcelTableAgent()
        
        print("="*60)
        print("Excel Table Agent - Demo")
        print("="*60)
        
        # Test data
        test_request = {
            'template': 'data_table',
            'variation': 'formatted_data_table',
            'data': {
                'headers': ['Facility', 'SO2 (mg/m³)', 'NOx (mg/m³)', 'PM10 (mg/m³)', 'Date'],
                'rows': [
                    ['Facility A', 45.2, 120.5, 15.3, '2025-01-15'],
                    ['Facility B', 38.7, 95.2, 12.8, '2025-01-15'],
                    ['Facility C', 52.1, 140.3, 18.9, '2025-01-15'],
                    ['Facility D', 41.5, 110.8, 14.2, '2025-01-15'],
                    ['Facility E', 48.9, 125.7, 16.5, '2025-01-15']
                ]
            },
            'output_format': 'excel',
            'filename': 'emission_data_test'
        }
        
        print("\n📊 Generating Excel table...")
        result = await agent.generate_table(test_request)
        
        if result['success']:
            print(f"✅ Success!")
            print(f"  File: {result['file_path']}")
            print(f"  Rows: {result['rows']}")
            print(f"  Columns: {result['columns']}")
        else:
            print(f"❌ Failed: {result['error']}")
        
        # Test CSV
        test_request['output_format'] = 'csv'
        test_request['filename'] = 'emission_data_csv_test'
        
        print("\n📄 Generating CSV file...")
        result = await agent.generate_table(test_request)
        
        if result['success']:
            print(f"✅ Success!")
            print(f"  File: {result['file_path']}")
        else:
            print(f"❌ Failed: {result['error']}")
        
        print("\n" + "="*60)
        print("✅ Demo complete!")
    
    asyncio.run(test_agent())
