"""
Vector Chart Agent - AI-gestützte Chart-Generierung

Dieser Agent nutzt LLM + RAG um automatisch Charts zu erstellen:
1. Intent Detection: Was möchte der User visualisieren?
2. Data Extraction: Daten aus Datenbank/RAG extrahieren
3. Chart Generation: Matplotlib/Plotly Chart erstellen
4. Export: PNG, SVG, PDF, PPTX
"""

import json
import logging
import base64
from typing import Dict, Any, List, Optional
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use('Agg')  # Non-GUI backend
import matplotlib.pyplot as plt
import seaborn as sns

logger = logging.getLogger(__name__)


class VectorChartAgent:
    """
    AI Agent für automatische Chart-Generierung
    
    Unterstützte Chart-Typen:
    - bar: Balkendiagramm
    - line: Liniendiagramm
    - pie: Kreisdiagramm
    - scatter: Streudiagramm
    - heatmap: Wärmekarte
    - table: Formatierte Tabelle
    """
    
    def __init__(
        self, 
        llm_service=None,
        rag_service=None,
        output_dir: str = "/tmp/veritas_charts"
    ):
        """
        Args:
            llm_service: LLM-Service für Intent Detection
            rag_service: RAG-Service für Datenextraktion
            output_dir: Verzeichnis für Chart-Exports
        """
        self.llm_service = llm_service
        self.rag_service = rag_service
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Seaborn-Theme für schönere Charts
        sns.set_theme(style="whitegrid")
        
        self.chart_templates = self._load_templates()
        logger.info("VectorChartAgent initialisiert")
    
    async def generate_chart(
        self, 
        user_prompt: str, 
        context: Optional[Dict[str, Any]] = None,
        template: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Hauptmethode: Chart-Generierung aus Nutzer-Prompt
        
        Args:
            user_prompt: "Erstelle ein Bar Chart mit BImSchG-Anlagen pro Kategorie"
            context: Optional zusätzlicher Kontext
            template: Optional Vorlagen-Name
            
        Returns:
            {
                'chart_type': 'bar',
                'title': 'Chart-Titel',
                'data': {...},
                'image_base64': 'iVBORw0KGgoAAAANS...',
                'exports': {
                    'png': '/tmp/chart_123.png',
                    'svg': '/tmp/chart_123.svg',
                    'pdf': '/tmp/chart_123.pdf',
                    'pptx': '/tmp/chart_123.pptx'
                }
            }
        """
        try:
            logger.info(f"Chart-Generierung gestartet: {user_prompt[:100]}")
            
            # Template verwenden wenn angegeben
            if template and template in self.chart_templates:
                intent = self.chart_templates[template].copy()
                logger.info(f"Template verwendet: {template}")
            else:
                # Step 1: Intent Detection via LLM
                intent = await self._detect_intent(user_prompt)
                logger.info(f"Intent erkannt: {intent.get('chart_type', 'unknown')}")
            
            # Step 2: Data Extraction
            data = await self._extract_data(intent, context or {})
            logger.info(f"Daten extrahiert: {len(data.get('values', []))} Datenpunkte")
            
            # Step 3: Chart Generation
            fig = self._generate_chart(intent['chart_type'], data, intent)
            logger.info("Chart generiert")
            
            # Step 4: Export-Vorbereitung
            exports = self._prepare_exports(fig, intent)
            logger.info(f"Exports vorbereitet: {list(exports.keys())}")
            
            plt.close(fig)  # Speicher freigeben
            
            return {
                'success': True,
                'chart_type': intent['chart_type'],
                'title': data.get('title', 'Chart'),
                'data': {
                    'labels': data.get('labels', []),
                    'values': data.get('values', [])
                },
                'image_base64': exports['png_base64'],
                'exports': {
                    'png': str(exports.get('png_path', '')),
                    'svg': str(exports.get('svg_path', '')),
                    'pdf': str(exports.get('pdf_path', '')),
                    'pptx': str(exports.get('pptx_path', ''))
                }
            }
        
        except Exception as e:
            logger.error(f"Chart-Generierung fehlgeschlagen: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'chart_type': 'error'
            }
    
    async def _detect_intent(self, user_prompt: str) -> Dict[str, Any]:
        """
        Nutze LLM um Chart-Intent zu erkennen
        
        Returns:
            {
                'chart_type': 'bar',
                'data_source': 'database' | 'rag' | 'example',
                'query': 'SELECT ...',
                'title': 'Chart-Titel',
                'x_label': 'X-Achse',
                'y_label': 'Y-Achse'
            }
        """
        if not self.llm_service:
            logger.warning("LLM-Service nicht verfügbar, nutze Fallback")
            return self._fallback_intent(user_prompt)
        
        system_prompt = """Du bist ein Chart-Spezifikations-Experte.
Analysiere die Nutzer-Anfrage und extrahiere Chart-Parameter im JSON-Format.

Unterstützte Chart-Typen: bar, line, pie, scatter, heatmap, table

Beispiel:
User: "Erstelle ein Bar Chart mit BImSchG-Anlagen pro 4. BImSchV-Kategorie"
Response: {
  "chart_type": "bar",
  "data_source": "database",
  "query": "SELECT nr_4bv, COUNT(*) as count FROM BImSchG GROUP BY nr_4bv ORDER BY count DESC LIMIT 10",
  "title": "BImSchG-Anlagen nach Kategorie",
  "x_label": "Kategorie (4. BImSchV)",
  "y_label": "Anzahl Anlagen"
}

Antworte NUR mit dem JSON-Objekt, keine zusätzlichen Erklärungen."""

        try:
            response = await self.llm_service.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.1,  # Präzise Extraktion
                max_tokens=500
            )
            
            # Parse JSON response
            response_text = response.get('text', '{}')
            # Extrahiere JSON aus Markdown-Code-Blöcken falls vorhanden
            if '```json' in response_text:
                response_text = response_text.split('```json')[1].split('```')[0].strip()
            elif '```' in response_text:
                response_text = response_text.split('```')[1].split('```')[0].strip()
            
            intent = json.loads(response_text)
            return intent
        
        except json.JSONDecodeError as e:
            logger.warning(f"JSON-Parse-Fehler: {e}, nutze Fallback")
            return self._fallback_intent(user_prompt)
        except Exception as e:
            logger.error(f"Intent-Detection fehlgeschlagen: {e}")
            return self._fallback_intent(user_prompt)
    
    def _fallback_intent(self, user_prompt: str) -> Dict[str, Any]:
        """Einfacher Fallback wenn LLM nicht verfügbar"""
        prompt_lower = user_prompt.lower()
        
        # Einfache Keyword-Erkennung
        if 'bar' in prompt_lower or 'balken' in prompt_lower:
            chart_type = 'bar'
        elif 'line' in prompt_lower or 'linien' in prompt_lower or 'zeit' in prompt_lower:
            chart_type = 'line'
        elif 'pie' in prompt_lower or 'kreis' in prompt_lower:
            chart_type = 'pie'
        elif 'scatter' in prompt_lower or 'streu' in prompt_lower:
            chart_type = 'scatter'
        elif 'heatmap' in prompt_lower or 'wärme' in prompt_lower:
            chart_type = 'heatmap'
        else:
            chart_type = 'bar'  # Default
        
        return {
            'chart_type': chart_type,
            'data_source': 'example',
            'title': 'Beispiel-Chart',
            'x_label': 'Kategorie',
            'y_label': 'Wert'
        }
    
    async def _extract_data(
        self, 
        intent: Dict[str, Any], 
        context: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Daten aus Datenbank/RAG extrahieren
        
        Strategien:
        1. SQL Query (wenn intent['data_source'] == 'database')
        2. RAG Search (wenn intent['data_source'] == 'rag')
        3. Example Data (wenn intent['data_source'] == 'example')
        """
        data_source = intent.get('data_source', 'example')
        
        if data_source == 'database':
            # TODO: SQL-Query ausführen (UDS3-Integration)
            logger.info("Database-Extraktion noch nicht implementiert, nutze Beispieldaten")
            data = self._get_example_data(intent['chart_type'])
        
        elif data_source == 'rag':
            # TODO: RAG-Suche
            logger.info("RAG-Extraktion noch nicht implementiert, nutze Beispieldaten")
            data = self._get_example_data(intent['chart_type'])
        
        else:
            # Example Data
            data = self._get_example_data(intent['chart_type'])
        
        # Intent-Felder zu Data hinzufügen
        data['title'] = intent.get('title', data.get('title', 'Chart'))
        data['x_label'] = intent.get('x_label', data.get('x_label', 'X'))
        data['y_label'] = intent.get('y_label', data.get('y_label', 'Y'))
        
        return data
    
    def _generate_chart(
        self, 
        chart_type: str, 
        data: Dict[str, Any],
        intent: Dict[str, Any]
    ) -> plt.Figure:
        """
        Matplotlib-Chart generieren
        
        Args:
            chart_type: 'bar', 'line', 'pie', 'scatter', 'heatmap'
            data: Chart-Daten
            intent: Intent-Informationen
            
        Returns:
            matplotlib.figure.Figure
        """
        fig, ax = plt.subplots(figsize=(10, 6))
        
        if chart_type == 'bar':
            labels = data.get('labels', [])
            values = data.get('values', [])
            
            bars = ax.bar(labels, values, color='steelblue', edgecolor='black', linewidth=1.2)
            ax.set_xlabel(data.get('x_label', 'Kategorie'), fontsize=11)
            ax.set_ylabel(data.get('y_label', 'Wert'), fontsize=11)
            
            # Werte auf Balken anzeigen
            for bar in bars:
                height = bar.get_height()
                ax.text(
                    bar.get_x() + bar.get_width() / 2.0,
                    height,
                    f'{int(height)}',
                    ha='center',
                    va='bottom',
                    fontsize=9
                )
        
        elif chart_type == 'line':
            x = data.get('x', list(range(len(data.get('y', [])))))
            y = data.get('y', [])
            
            ax.plot(x, y, marker='o', linewidth=2, markersize=6, color='steelblue')
            ax.set_xlabel(data.get('x_label', 'X'), fontsize=11)
            ax.set_ylabel(data.get('y_label', 'Y'), fontsize=11)
            ax.grid(True, alpha=0.3)
        
        elif chart_type == 'pie':
            labels = data.get('labels', [])
            values = data.get('values', [])
            
            colors = sns.color_palette('pastel', len(labels))
            wedges, texts, autotexts = ax.pie(
                values,
                labels=labels,
                autopct='%1.1f%%',
                startangle=90,
                colors=colors,
                wedgeprops=dict(edgecolor='white', linewidth=2)
            )
            
            # Text-Formatierung
            for text in texts:
                text.set_fontsize(10)
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(9)
                autotext.set_fontweight('bold')
        
        elif chart_type == 'scatter':
            x = data.get('x', [])
            y = data.get('y', [])
            
            ax.scatter(x, y, s=100, alpha=0.6, color='steelblue', edgecolors='black', linewidth=1)
            ax.set_xlabel(data.get('x_label', 'X'), fontsize=11)
            ax.set_ylabel(data.get('y_label', 'Y'), fontsize=11)
            ax.grid(True, alpha=0.3)
        
        elif chart_type == 'heatmap':
            matrix = data.get('matrix', [[1, 2], [3, 4]])
            
            sns.heatmap(
                matrix,
                annot=True,
                fmt='.0f',
                cmap='YlOrRd',
                ax=ax,
                cbar_kws={'label': 'Wert'}
            )
            ax.set_xlabel(data.get('x_label', 'X'), fontsize=11)
            ax.set_ylabel(data.get('y_label', 'Y'), fontsize=11)
        
        else:
            # Fallback: Bar Chart
            ax.bar(['A', 'B', 'C'], [1, 2, 3])
        
        # Titel
        ax.set_title(data.get('title', 'Chart'), fontsize=14, fontweight='bold', pad=20)
        
        plt.tight_layout()
        
        return fig
    
    def _prepare_exports(
        self, 
        fig: plt.Figure, 
        intent: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Exports vorbereiten: PNG, SVG, PDF
        
        Returns:
            {
                'png_base64': 'iVBORw0KGgo...',
                'png_path': Path('/tmp/chart_123.png'),
                'svg_path': Path('/tmp/chart_123.svg'),
                'pdf_path': Path('/tmp/chart_123.pdf')
            }
        """
        import time
        
        # Eindeutiger Dateiname
        timestamp = int(time.time() * 1000)
        title_slug = intent.get('title', 'chart').replace(' ', '_').lower()[:30]
        base_name = f"{title_slug}_{timestamp}"
        
        exports = {}
        
        # PNG (Base64 für Frontend)
        png_buf = BytesIO()
        fig.savefig(png_buf, format='png', dpi=150, bbox_inches='tight')
        png_buf.seek(0)
        png_base64 = base64.b64encode(png_buf.read()).decode('utf-8')
        exports['png_base64'] = png_base64
        
        # PNG (Datei)
        png_path = self.output_dir / f"{base_name}.png"
        fig.savefig(png_path, format='png', dpi=150, bbox_inches='tight')
        exports['png_path'] = png_path
        
        # SVG
        svg_path = self.output_dir / f"{base_name}.svg"
        fig.savefig(svg_path, format='svg', bbox_inches='tight')
        exports['svg_path'] = svg_path
        
        # PDF
        pdf_path = self.output_dir / f"{base_name}.pdf"
        fig.savefig(pdf_path, format='pdf', bbox_inches='tight')
        exports['pdf_path'] = pdf_path
        
        # PPTX (optional, erfordert python-pptx)
        try:
            pptx_path = self._create_pptx(fig, intent, base_name)
            exports['pptx_path'] = pptx_path
        except Exception as e:
            logger.warning(f"PPTX-Export fehlgeschlagen: {e}")
        
        return exports
    
    def _create_pptx(
        self, 
        fig: plt.Figure, 
        intent: Dict[str, Any],
        base_name: str
    ) -> Path:
        """
        PowerPoint-Präsentation erstellen
        
        Erfordert: pip install python-pptx
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.warning("python-pptx nicht installiert, PPTX-Export übersprungen")
            return None
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Chart als Bild einfügen
        img_buf = BytesIO()
        fig.savefig(img_buf, format='png', dpi=150, bbox_inches='tight')
        img_buf.seek(0)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        slide.shapes.add_picture(img_buf, left, top, width=width)
        
        # Titel hinzufügen
        title_box = slide.shapes.title
        if title_box:
            title_box.text = intent.get('title', 'Chart')
        
        pptx_path = self.output_dir / f"{base_name}.pptx"
        prs.save(str(pptx_path))
        
        return pptx_path
    
    def _load_templates(self) -> Dict[str, Dict[str, Any]]:
        """
        Vordefinierte Chart-Templates laden
        
        Returns:
            {
                'template_name': {
                    'chart_type': 'bar',
                    'data_source': 'database',
                    'query': 'SELECT ...',
                    'title': '...',
                    'x_label': '...',
                    'y_label': '...'
                }
            }
        """
        return {
            'bimschg_overview': {
                'chart_type': 'bar',
                'data_source': 'database',
                'query': 'SELECT nr_4bv, COUNT(*) as count FROM BImSchG GROUP BY nr_4bv ORDER BY count DESC LIMIT 10',
                'title': 'BImSchG-Anlagen nach Kategorie',
                'x_label': '4. BImSchV Nummer',
                'y_label': 'Anzahl Anlagen'
            },
            'wka_leistung': {
                'chart_type': 'pie',
                'data_source': 'database',
                'query': 'SELECT status, SUM(leistung) as total FROM WKA GROUP BY status',
                'title': 'WKA-Leistung nach Status'
            },
            'anlagenverteilung': {
                'chart_type': 'pie',
                'data_source': 'database',
                'query': 'SELECT anlart_4bv, COUNT(*) as count FROM BImSchG GROUP BY anlart_4bv LIMIT 8',
                'title': 'Verteilung der Anlagentypen'
            },
            'zeitreihe_genehmigungen': {
                'chart_type': 'line',
                'data_source': 'database',
                'query': "SELECT YEAR(datum) as year, COUNT(*) as count FROM BImSchG WHERE YEAR(datum) BETWEEN 2010 AND 2024 GROUP BY YEAR(datum) ORDER BY year",
                'title': 'Genehmigungen pro Jahr (2010-2024)',
                'x_label': 'Jahr',
                'y_label': 'Anzahl Genehmigungen'
            }
        }
    
    def _get_example_data(self, chart_type: str) -> Dict[str, Any]:
        """
        Beispiel-Daten für Demo-Charts
        
        Returns:
            {
                'labels': [...],
                'values': [...],
                'title': '...',
                'x_label': '...',
                'y_label': '...'
            }
        """
        examples = {
            'bar': {
                'labels': ['1.1 Feuerung', '1.2 Gasturbine', '7.1 Tierhaltung', '8.1 Abfall', '10.1 Sonstige'],
                'values': [850, 520, 1200, 340, 650],
                'title': 'BImSchG-Anlagen nach Kategorie (Beispieldaten)',
                'x_label': '4. BImSchV Kategorie',
                'y_label': 'Anzahl Anlagen'
            },
            'line': {
                'x': [2015, 2016, 2017, 2018, 2019, 2020, 2021, 2022, 2023, 2024],
                'y': [120, 145, 160, 155, 180, 195, 210, 225, 240, 260],
                'title': 'Genehmigungen pro Jahr (Beispieldaten)',
                'x_label': 'Jahr',
                'y_label': 'Anzahl Genehmigungen'
            },
            'pie': {
                'labels': ['In Betrieb', 'Im Genehmigungsverfahren', 'Stillgelegt', 'Im Bau'],
                'values': [3500, 850, 420, 230],
                'title': 'WKA-Status-Verteilung (Beispieldaten)'
            },
            'scatter': {
                'x': [10, 20, 30, 40, 50, 60, 70, 80, 90, 100],
                'y': [25, 35, 28, 48, 55, 62, 58, 75, 82, 88],
                'title': 'Leistung vs. Nabenhöhe (Beispieldaten)',
                'x_label': 'Nabenhöhe (m)',
                'y_label': 'Leistung (MW)'
            },
            'heatmap': {
                'matrix': [
                    [10, 20, 30, 15],
                    [25, 35, 20, 40],
                    [15, 25, 45, 30],
                    [30, 15, 25, 35]
                ],
                'title': 'Anlagen-Heatmap (Beispieldaten)',
                'x_label': 'Region',
                'y_label': 'Kategorie'
            }
        }
        
        return examples.get(chart_type, examples['bar'])
    
    def list_templates(self) -> List[Dict[str, str]]:
        """
        Liste aller verfügbaren Templates
        
        Returns:
            [
                {'name': 'bimschg_overview', 'title': 'BImSchG-Anlagen nach Kategorie', 'type': 'bar'},
                ...
            ]
        """
        return [
            {
                'name': template_name,
                'title': template.get('title', template_name),
                'type': template.get('chart_type', 'bar')
            }
            for template_name, template in self.chart_templates.items()
        ]


# Standalone-Test
if __name__ == '__main__':
    import asyncio
    
    async def test():
        agent = VectorChartAgent()
        
        # Test 1: Bar Chart
        result = await agent.generate_chart(
            "Erstelle ein Bar Chart mit BImSchG-Anlagen",
            template='bimschg_overview'
        )
        print(f"✅ Bar Chart: {result['success']}")
        print(f"   PNG: {result['exports']['png']}")
        
        # Test 2: Pie Chart
        result = await agent.generate_chart(
            "Zeige ein Pie Chart der WKA-Leistung",
            template='wka_leistung'
        )
        print(f"✅ Pie Chart: {result['success']}")
        
        # Test 3: Line Chart
        result = await agent.generate_chart(
            "Liniendiagramm: Genehmigungen pro Jahr",
            template='zeitreihe_genehmigungen'
        )
        print(f"✅ Line Chart: {result['success']}")
        
        # Test 4: Templates auflisten
        templates = agent.list_templates()
        print(f"\n📋 Verfügbare Templates: {len(templates)}")
        for t in templates:
            print(f"   - {t['name']}: {t['title']} ({t['type']})")
    
    asyncio.run(test())
