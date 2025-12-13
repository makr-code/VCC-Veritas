"""
Presentation Canvas Agent - Erweiterte Präsentations-Generierung

Dieser Agent erweitert den Vector Chart Agent um:
1. Bildbeschreibende Sprache (Visual Description Language - VDL)
2. Canvas-basierte Präsentationserstellung
3. Integration mit LLM für visuelle Beschreibungen
4. Vorbereitung für AI-Bildgenerator-Integration

VDL (Visual Description Language) Format:
{
    "slides": [
        {
            "layout": "title_slide" | "content" | "two_column" | "chart" | "image",
            "elements": [
                {
                    "type": "text" | "shape" | "chart" | "image" | "icon",
                    "position": {"x": int, "y": int},
                    "size": {"width": int, "height": int},
                    "properties": {...}
                }
            ]
        }
    ]
}
"""

import json
import logging
import base64
import time
import math
from typing import Dict, Any, List, Optional, Tuple
from io import BytesIO
from pathlib import Path

# PIL für erweiterte Grafik-Operationen
from PIL import Image, ImageDraw, ImageFont

# Tkinter für Canvas (optional, nur für interaktive UI)
try:
    import tkinter as tk
    from tkinter import Canvas, font as tkfont
    from PIL import ImageTk
    TKINTER_AVAILABLE = True
except ImportError:
    TKINTER_AVAILABLE = False
    tk = None
    ImageTk = None

# PowerPoint Integration
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class VisualDescriptionLanguage:
    """
    VDL (Visual Description Language) Parser und Validator
    
    Definiert eine strukturierte Sprache zur Beschreibung visueller Elemente,
    die vom LLM erzeugt und vom Canvas Agent interpretiert werden kann.
    """
    
    ELEMENT_TYPES = [
        'text', 'shape', 'chart', 'image', 'icon', 'line', 'arrow',
        'connector', 'flowchart', 'org_chart', 'process_flow', 
        'cycle_diagram', 'pyramid'
    ]
    LAYOUT_TYPES = ['title_slide', 'content', 'two_column', 'chart', 'image', 'blank']
    SHAPE_TYPES = [
        # Basis-Formen
        'rectangle', 'rounded_rectangle', 'oval', 'triangle',
        'diamond', 'pentagon', 'hexagon', 'octagon',
        
        # Pfeile (häufig genutzte)
        'right_arrow', 'left_arrow', 'up_arrow', 'down_arrow',
        'left_right_arrow', 'circular_arrow', 'bent_arrow',
        
        # Flussdiagramm-Formen
        'flowchart_process', 'flowchart_decision', 'flowchart_terminator',
        'flowchart_data', 'flowchart_document', 'flowchart_database',  # Note: database = magnetic_disk
        
        # Callouts
        'cloud_callout', 'oval_callout',
        
        # Sterne
        'star_5_point', 'star_6_point'
    ]
    CONNECTOR_TYPES = ['straight', 'elbow', 'curve']
    
    @classmethod
    def validate(cls, vdl_spec: Dict[str, Any]) -> Tuple[bool, Optional[str]]:
        """
        Validiert eine VDL-Spezifikation
        
        Returns:
            (is_valid, error_message)
        """
        if 'slides' not in vdl_spec:
            return False, "Missing 'slides' key"
        
        if not isinstance(vdl_spec['slides'], list):
            return False, "'slides' must be a list"
        
        for i, slide in enumerate(vdl_spec['slides']):
            if 'layout' not in slide:
                return False, f"Slide {i}: Missing 'layout'"
            
            if slide['layout'] not in cls.LAYOUT_TYPES:
                return False, f"Slide {i}: Invalid layout '{slide['layout']}'"
            
            if 'elements' in slide:
                for j, element in enumerate(slide['elements']):
                    if 'type' not in element:
                        return False, f"Slide {i}, Element {j}: Missing 'type'"
                    
                    if element['type'] not in cls.ELEMENT_TYPES:
                        return False, f"Slide {i}, Element {j}: Invalid type '{element['type']}'"
        
        return True, None
    
    @classmethod
    def create_example(cls) -> Dict[str, Any]:
        """Erstellt ein Beispiel-VDL für Demonstrations-Zwecke"""
        return {
            "metadata": {
                "title": "Beispiel-Präsentation",
                "author": "VERITAS Canvas Agent",
                "theme": "professional"
            },
            "slides": [
                {
                    "layout": "title_slide",
                    "elements": [
                        {
                            "type": "text",
                            "content": "Präsentationstitel",
                            "position": {"x": 50, "y": 200},
                            "size": {"width": 700, "height": 100},
                            "properties": {
                                "font_size": 44,
                                "font_weight": "bold",
                                "align": "center",
                                "color": "#1f4788"
                            }
                        },
                        {
                            "type": "text",
                            "content": "Untertitel oder Autor",
                            "position": {"x": 50, "y": 320},
                            "size": {"width": 700, "height": 50},
                            "properties": {
                                "font_size": 24,
                                "align": "center",
                                "color": "#666666"
                            }
                        }
                    ]
                },
                {
                    "layout": "content",
                    "elements": [
                        {
                            "type": "text",
                            "content": "Folientitel",
                            "position": {"x": 50, "y": 50},
                            "size": {"width": 700, "height": 60},
                            "properties": {
                                "font_size": 32,
                                "font_weight": "bold",
                                "color": "#1f4788"
                            }
                        },
                        {
                            "type": "shape",
                            "shape": "rectangle",
                            "position": {"x": 100, "y": 150},
                            "size": {"width": 200, "height": 150},
                            "properties": {
                                "fill_color": "#e8f4f8",
                                "border_color": "#1f4788",
                                "border_width": 2
                            }
                        },
                        {
                            "type": "text",
                            "content": "Inhaltselement",
                            "position": {"x": 120, "y": 200},
                            "size": {"width": 160, "height": 50},
                            "properties": {
                                "font_size": 16,
                                "align": "center"
                            }
                        }
                    ]
                }
            ]
        }


class PresentationCanvasAgent:
    """
    Canvas Agent für erweiterte Präsentations-Generierung
    
    Funktionen:
    1. VDL-basierte Präsentationserstellung
    2. LLM-Integration für visuelle Beschreibungen
    3. Canvas-Rendering (Tkinter)
    4. PowerPoint-Export
    5. Vorbereitung für AI-Bildgenerator
    """
    
    def __init__(
        self,
        llm_service=None,
        output_dir: str = "/tmp/veritas_presentations"
    ):
        """
        Args:
            llm_service: LLM-Service für VDL-Generierung
            output_dir: Verzeichnis für Präsentations-Exports
        """
        self.llm_service = llm_service
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info("PresentationCanvasAgent initialisiert")
    
    @staticmethod
    def _get_font(size: int = 16) -> ImageFont.FreeTypeFont:
        """
        Cross-platform font loading mit Fallback
        
        Versucht verschiedene Font-Pfade für Linux, Windows, macOS
        """
        font_paths = [
            # Linux
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/TTF/DejaVuSans.ttf",
            # macOS
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
            # Windows
            "C:\\Windows\\Fonts\\arial.ttf",
            "C:\\Windows\\Fonts\\calibri.ttf",
        ]
        
        for font_path in font_paths:
            try:
                return ImageFont.truetype(font_path, size)
            except (OSError, IOError):
                continue
        
        # Fallback: Default-Font
        logger.warning("Kein TrueType-Font gefunden, nutze Default-Font")
        return ImageFont.load_default()
    
    async def generate_presentation(
        self,
        user_prompt: str,
        context: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Präsentation aus Nutzer-Prompt generieren
        
        Workflow:
        1. LLM generiert VDL-Spezifikation
        2. VDL wird validiert
        3. Canvas-Rendering
        4. Export zu PPTX/PNG
        
        Args:
            user_prompt: "Erstelle eine Präsentation über BImSchG-Anlagen mit 3 Folien"
            context: Optional zusätzlicher Kontext
            
        Returns:
            {
                'success': bool,
                'vdl': dict,
                'slides': [{'image_base64': str, 'png_path': str}, ...],
                'pptx_path': str
            }
        """
        try:
            logger.info(f"Präsentations-Generierung gestartet: {user_prompt[:100]}")
            
            # Step 1: VDL via LLM generieren
            vdl = await self._generate_vdl(user_prompt, context or {})
            logger.info(f"VDL generiert: {len(vdl.get('slides', []))} Folien")
            
            # Step 2: VDL validieren
            is_valid, error = VisualDescriptionLanguage.validate(vdl)
            if not is_valid:
                raise ValueError(f"Ungültige VDL: {error}")
            
            # Step 3: Slides rendern
            slides = await self._render_slides(vdl)
            logger.info(f"{len(slides)} Folien gerendert")
            
            # Step 4: PPTX erstellen
            pptx_path = None
            if PPTX_AVAILABLE:
                pptx_path = self._create_pptx(vdl, slides)
                logger.info(f"PPTX erstellt: {pptx_path}")
            
            return {
                'success': True,
                'vdl': vdl,
                'slides': slides,
                'pptx_path': str(pptx_path) if pptx_path else None,
                'slide_count': len(slides)
            }
        
        except Exception as e:
            logger.error(f"Präsentations-Generierung fehlgeschlagen: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
    
    async def _generate_vdl(
        self,
        user_prompt: str,
        context: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        VDL-Spezifikation via LLM generieren
        
        Das LLM wird instruiert, eine strukturierte visuelle Beschreibung
        zu erstellen, die alle Elemente einer Präsentation definiert.
        """
        if not self.llm_service:
            logger.warning("LLM-Service nicht verfügbar, nutze Beispiel-VDL")
            return VisualDescriptionLanguage.create_example()
        
        system_prompt = """Du bist ein Experte für visuelle Präsentationsgestaltung.
Erstelle eine Visual Description Language (VDL) Spezifikation basierend auf der Nutzer-Anfrage.

VDL-Format (JSON):
{
  "metadata": {
    "title": "Präsentationstitel",
    "author": "Autor",
    "theme": "professional" | "minimal" | "colorful"
  },
  "use_native_shapes": true,  // Nutze native PowerPoint-Shapes (editierbar)
  "slides": [
    {
      "layout": "title_slide" | "content" | "two_column" | "chart" | "image",
      "elements": [
        {
          "type": "text" | "shape" | "connector" | "flowchart" | "org_chart" | "cycle_diagram" | "chart" | "image",
          "content": "Text oder Beschreibung",
          "position": {"x": int, "y": int},
          "size": {"width": int, "height": int},
          "properties": {
            "font_size": int,
            "color": "#hexcolor",
            "align": "left" | "center" | "right",
            ...
          }
        }
      ]
    }
  ]
}

Elemente-Typen:
- text: Textelemente (Titel, Inhalt, Listen)
- shape: Formen (rectangle, rounded_rectangle, oval, triangle, diamond, pentagon, hexagon, octagon)
- connector: Verbindungslinien (straight, elbow, curve)
- flowchart: Flussdiagramm-Template mit steps-Array
- org_chart: Organigramm-Template mit levels-Array
- cycle_diagram: Zyklisches Diagramm-Template mit steps-Array
- chart: Verweise auf Chart-Agent (chart_id oder chart_spec)
- image: Platzhalter für Bilder (ai_generated oder file)

Shape-Typen (häufig genutzt):
- Basis: rectangle, rounded_rectangle, oval, triangle, diamond, pentagon, hexagon, octagon
- Pfeile: right_arrow, left_arrow, up_arrow, down_arrow, left_right_arrow, circular_arrow, bent_arrow
- Flowchart: flowchart_process, flowchart_decision, flowchart_terminator, flowchart_data, flowchart_document, flowchart_database
- Callouts: cloud_callout, oval_callout
- Sterne: star_5_point, star_6_point

Diagram-Templates (verwende diese für komplexe Strukturen):
1. flowchart:
   {
     "type": "flowchart",
     "steps": [
       {"shape": "flowchart_terminator", "text": "Start", "color": "#70ad47"},
       {"shape": "flowchart_process", "text": "Schritt 1", "color": "#4472c4"},
       {"shape": "flowchart_decision", "text": "Entscheidung?", "color": "#ed7d31"},
       {"shape": "flowchart_terminator", "text": "Ende", "color": "#70ad47"}
     ],
     "position": {"x": 100, "y": 100}
   }

2. org_chart:
   {
     "type": "org_chart",
     "levels": [
       ["CEO"],
       ["Manager 1", "Manager 2"],
       ["Team A", "Team B", "Team C"]
     ],
     "position": {"x": 100, "y": 100}
   }

3. cycle_diagram:
   {
     "type": "cycle_diagram",
     "steps": ["Plan", "Do", "Check", "Act"],
     "position": {"x": 400, "y": 300}
   }

Koordinaten-System:
- Canvas: 800x600 (Standard)
- x: 0-800 (links nach rechts)
- y: 0-600 (oben nach unten)

Best Practices:
- Nutze Templates (flowchart, org_chart, cycle_diagram) statt manuelle Shape-Positionierung
- Verwende konsistente Farben (#4472c4 für primary, #70ad47 für success, #ed7d31 für warnings)
- Lasse Platz zwischen Elementen (mind. 20px)
- Zentriere wichtige Elemente
- Nutze native_shapes für editierbare PowerPoint-Präsentationen

Antworte NUR mit dem VDL-JSON, keine Erklärungen."""

        try:
            response = await self.llm_service.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.2,  # Kreativ aber strukturiert
                max_tokens=2000
            )
            
            # Parse JSON response
            response_text = response.get('text', '{}')
            
            # Extrahiere JSON aus Markdown-Code-Blöcken
            if '```json' in response_text:
                response_text = response_text.split('```json')[1].split('```')[0].strip()
            elif '```' in response_text:
                response_text = response_text.split('```')[1].split('```')[0].strip()
            
            vdl = json.loads(response_text)
            return vdl
        
        except Exception as e:
            logger.error(f"VDL-Generierung fehlgeschlagen: {e}")
            return VisualDescriptionLanguage.create_example()
    
    async def _render_slides(self, vdl: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        VDL zu Canvas-Bildern rendern
        
        Returns:
            List of {'image_base64': str, 'png_path': str}
        """
        slides = []
        
        for i, slide_spec in enumerate(vdl.get('slides', [])):
            # Canvas erstellen (PIL-basiert für Headless-Rendering)
            img = Image.new('RGB', (800, 600), 'white')
            draw = ImageDraw.Draw(img)
            
            # Layout-abhängige Hintergrundgestaltung
            self._apply_layout(img, draw, slide_spec.get('layout', 'blank'))
            
            # Elemente rendern
            for element in slide_spec.get('elements', []):
                self._render_element(img, draw, element)
            
            # Speichern
            timestamp = int(time.time() * 1000)
            png_path = self.output_dir / f"slide_{i+1}_{timestamp}.png"
            img.save(png_path, 'PNG', dpi=(150, 150))
            
            # Base64 encodieren
            buf = BytesIO()
            img.save(buf, format='PNG')
            buf.seek(0)
            image_base64 = base64.b64encode(buf.read()).decode('utf-8')
            
            slides.append({
                'image_base64': image_base64,
                'png_path': str(png_path),
                'slide_number': i + 1
            })
        
        return slides
    
    def _apply_layout(self, img: Image.Image, draw: ImageDraw.Draw, layout: str):
        """Layout-spezifisches Styling anwenden"""
        if layout == 'title_slide':
            # Gradient-ähnlicher Effekt (vereinfacht)
            for y in range(0, 200):
                alpha = int((y / 200) * 30)
                color = (31 + alpha, 71 + alpha, 136 + alpha)
                draw.rectangle([(0, y), (800, y+1)], fill=color)
        
        elif layout in ['content', 'two_column']:
            # Subtile Kopfzeile
            draw.rectangle([(0, 0), (800, 80)], fill='#f8f9fa')
            draw.line([(0, 80), (800, 80)], fill='#1f4788', width=3)
    
    def _render_element(self, img: Image.Image, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Einzelnes VDL-Element rendern"""
        element_type = element.get('type')
        pos = element.get('position', {'x': 0, 'y': 0})
        size = element.get('size', {'width': 100, 'height': 50})
        props = element.get('properties', {})
        
        if element_type == 'text':
            self._render_text(draw, element, pos, size, props)
        
        elif element_type == 'shape':
            self._render_shape(draw, element, pos, size, props)
        
        elif element_type == 'connector':
            self._render_connector(draw, element)
        
        elif element_type == 'flowchart':
            self._render_flowchart_template(draw, element)
        
        elif element_type == 'org_chart':
            self._render_org_chart_template(draw, element)
        
        elif element_type == 'process_flow':
            self._render_process_flow_template(draw, element)
        
        elif element_type == 'cycle_diagram':
            self._render_cycle_diagram_template(draw, element)
        
        elif element_type == 'chart':
            # Platzhalter für Chart (wird via Chart-Agent erstellt)
            self._render_chart_placeholder(draw, pos, size, element)
        
        elif element_type == 'image':
            # Platzhalter für AI-generiertes Bild
            self._render_image_placeholder(draw, pos, size, element)
    
    def _render_text(
        self,
        draw: ImageDraw.Draw,
        element: Dict[str, Any],
        pos: Dict[str, int],
        size: Dict[str, int],
        props: Dict[str, Any]
    ):
        """Text-Element rendern"""
        content = element.get('content', '')
        font_size = props.get('font_size', 16)
        color = props.get('color', '#000000')
        
        # Load font using cross-platform helper
        font = self._get_font(font_size)
        
        # Text-Position (vereinfacht, ohne Alignment)
        x = pos.get('x', 0)
        y = pos.get('y', 0)
        
        draw.text((x, y), content, fill=color, font=font)
    
    def _render_shape(
        self,
        draw: ImageDraw.Draw,
        element: Dict[str, Any],
        pos: Dict[str, int],
        size: Dict[str, int],
        props: Dict[str, Any]
    ):
        """Form-Element rendern"""
        shape_type = element.get('shape', 'rectangle')
        x, y = pos.get('x', 0), pos.get('y', 0)
        w, h = size.get('width', 100), size.get('height', 50)
        
        fill_color = props.get('fill_color', '#e8f4f8')
        border_color = props.get('border_color', '#1f4788')
        border_width = props.get('border_width', 2)
        
        if shape_type == 'rectangle':
            draw.rectangle(
                [(x, y), (x + w, y + h)],
                fill=fill_color,
                outline=border_color,
                width=border_width
            )
        
        elif shape_type == 'circle':
            draw.ellipse(
                [(x, y), (x + w, y + h)],
                fill=fill_color,
                outline=border_color,
                width=border_width
            )
    
    def _render_chart_placeholder(
        self,
        draw: ImageDraw.Draw,
        pos: Dict[str, int],
        size: Dict[str, int],
        element: Dict[str, Any]
    ):
        """Platzhalter für Chart (Integration mit VectorChartAgent)"""
        x, y = pos.get('x', 0), pos.get('y', 0)
        w, h = size.get('width', 400), size.get('height', 300)
        
        # Platzhalter-Box
        draw.rectangle([(x, y), (x + w, y + h)], fill='#f0f0f0', outline='#1f4788', width=2)
        
        # Text
        chart_spec = element.get('chart_spec', 'Chart')
        font = self._get_font(14)
        
        draw.text((x + 10, y + 10), f"[Chart: {chart_spec}]", fill='#666666', font=font)
        draw.text((x + 10, y + 30), "→ Chart Agent Integration", fill='#999999', font=font)
    
    def _render_image_placeholder(
        self,
        draw: ImageDraw.Draw,
        pos: Dict[str, int],
        size: Dict[str, int],
        element: Dict[str, Any]
    ):
        """Platzhalter für AI-generiertes Bild"""
        x, y = pos.get('x', 0), pos.get('y', 0)
        w, h = size.get('width', 400), size.get('height', 300)
        
        # Platzhalter-Box
        draw.rectangle([(x, y), (x + w, y + h)], fill='#fff8e1', outline='#ff9800', width=2)
        
        # Text
        ai_prompt = element.get('ai_prompt', 'AI Image')
        font = self._get_font(14)
        
        draw.text((x + 10, y + 10), f"[AI Image: {ai_prompt[:40]}]", fill='#ff9800', font=font)
        draw.text((x + 10, y + 30), "→ AI Image Generator", fill='#ffa726', font=font)
    
    def _render_connector(self, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Connector/Verbindungslinie rendern"""
        connector_type = element.get('connector_type', 'straight')
        start = element.get('start', {'x': 0, 'y': 0})
        end = element.get('end', {'x': 100, 'y': 100})
        props = element.get('properties', {})
        
        width = props.get('width', 2)
        color = props.get('color', '#000000')
        
        if connector_type == 'straight':
            draw.line(
                [(start['x'], start['y']), (end['x'], end['y'])],
                fill=color,
                width=width
            )
        elif connector_type == 'elbow':
            # L-förmige Verbindung (90° Winkel)
            mid_x = start['x']
            mid_y = end['y']
            draw.line(
                [(start['x'], start['y']), (mid_x, mid_y), (end['x'], end['y'])],
                fill=color,
                width=width
            )
    
    def _render_flowchart_template(self, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Flowchart-Diagramm aus Template rendern"""
        steps = element.get('steps', [])
        pos = element.get('position', {'x': 100, 'y': 100})
        
        y_spacing = 80
        x_pos = pos['x']
        start_y = pos['y']
        
        for idx, step in enumerate(steps):
            y_pos = start_y + idx * y_spacing
            
            # Shape rendern
            shape_element = {
                'type': 'shape',
                'shape': step.get('shape', 'rectangle'),
                'content': step.get('text', ''),
                'position': {'x': x_pos, 'y': y_pos},
                'size': {'width': 150, 'height': 60},
                'properties': {
                    'fill_color': step.get('color', '#4472c4'),
                    'border_color': '#000000',
                    'border_width': 2
                }
            }
            self._render_shape(draw, shape_element, 
                             shape_element['position'], 
                             shape_element['size'],
                             shape_element['properties'])
            
            # Connector zum nächsten Schritt
            if idx < len(steps) - 1:
                self._render_connector(draw, {
                    'connector_type': 'straight',
                    'start': {'x': x_pos + 75, 'y': y_pos + 60},
                    'end': {'x': x_pos + 75, 'y': y_pos + y_spacing},
                    'properties': {'width': 2, 'color': '#000000'}
                })
    
    def _render_org_chart_template(self, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Organigramm aus Template rendern"""
        levels = element.get('levels', [])
        pos = element.get('position', {'x': 100, 'y': 100})
        
        y_spacing = 100
        x_spacing = 120
        
        for level_idx, level in enumerate(levels):
            y_pos = pos['y'] + level_idx * y_spacing
            total_width = len(level) * x_spacing
            start_x = pos['x'] + (600 - total_width) / 2  # Zentriert in 600px Breite
            
            for node_idx, node_text in enumerate(level):
                x_pos = start_x + node_idx * x_spacing
                
                # Node rendern
                shape_element = {
                    'type': 'shape',
                    'shape': 'rounded_rectangle',
                    'content': node_text,
                    'position': {'x': x_pos, 'y': y_pos},
                    'size': {'width': 100, 'height': 50},
                    'properties': {
                        'fill_color': '#4472c4',
                        'border_color': '#000000',
                        'border_width': 2
                    }
                }
                self._render_shape(draw, shape_element,
                                 shape_element['position'],
                                 shape_element['size'],
                                 shape_element['properties'])
                
                # Connector zur vorherigen Ebene
                if level_idx > 0:
                    parent_idx = node_idx // 2
                    prev_level_count = len(levels[level_idx - 1])
                    prev_total_width = prev_level_count * x_spacing
                    prev_start_x = pos['x'] + (600 - prev_total_width) / 2
                    parent_x = prev_start_x + parent_idx * x_spacing
                    
                    self._render_connector(draw, {
                        'connector_type': 'elbow',
                        'start': {'x': parent_x + 50, 'y': y_pos - y_spacing + 50},
                        'end': {'x': x_pos + 50, 'y': y_pos},
                        'properties': {'width': 2, 'color': '#000000'}
                    })
    
    def _render_process_flow_template(self, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Prozessfluss-Diagramm rendern"""
        # Ähnlich wie flowchart_template, aber mit spezifischen Process-Shapes
        steps = element.get('steps', [])
        self._render_flowchart_template(draw, element)
    
    def _render_cycle_diagram_template(self, draw: ImageDraw.Draw, element: Dict[str, Any]):
        """Zyklisches Diagramm rendern"""
        import math
        
        steps = element.get('steps', [])
        pos = element.get('position', {'x': 400, 'y': 300})
        
        center_x = pos['x']
        center_y = pos['y']
        radius = 120
        n = len(steps)
        
        for idx, step_text in enumerate(steps):
            # Position auf dem Kreis berechnen
            angle = (2 * math.pi * idx / n) - (math.pi / 2)  # Start oben
            x = center_x + radius * math.cos(angle)
            y = center_y + radius * math.sin(angle)
            
            # Shape rendern
            shape_element = {
                'type': 'shape',
                'shape': 'oval',
                'content': step_text,
                'position': {'x': int(x - 40), 'y': int(y - 40)},
                'size': {'width': 80, 'height': 80},
                'properties': {
                    'fill_color': '#4472c4',
                    'border_color': '#000000',
                    'border_width': 2
                }
            }
            self._render_shape(draw, shape_element,
                             shape_element['position'],
                             shape_element['size'],
                             shape_element['properties'])
    
    def _create_pptx(
        self,
        vdl: Dict[str, Any],
        slides: List[Dict[str, Any]]
    ) -> Optional[Path]:
        """PowerPoint-Präsentation erstellen"""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx nicht verfügbar, PPTX-Export übersprungen")
            return None
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Metadata
            metadata = vdl.get('metadata', {})
            prs.core_properties.title = metadata.get('title', 'Präsentation')
            prs.core_properties.author = metadata.get('author', 'VERITAS Canvas Agent')
            
            # Option: Native Shapes vs. Bilder
            use_native_shapes = vdl.get('use_native_shapes', False)
            
            if use_native_shapes:
                # Native PowerPoint Shapes erstellen (direkt manipulierbar)
                for slide_spec in vdl.get('slides', []):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
                    self._add_native_shapes_to_slide(slide, slide_spec)
            else:
                # Slides als Bilder hinzufügen (Fallback)
                for slide_data in slides:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
                    
                    # Slide-Bild einfügen
                    png_path = slide_data['png_path']
                    left = Inches(0)
                    top = Inches(0)
                    slide.shapes.add_picture(png_path, left, top, width=Inches(10))
            
            # Speichern
            timestamp = int(time.time() * 1000)
            title_slug = metadata.get('title', 'presentation').replace(' ', '_').lower()[:30]
            pptx_path = self.output_dir / f"{title_slug}_{timestamp}.pptx"
            prs.save(str(pptx_path))
            
            return pptx_path
        
        except Exception as e:
            logger.error(f"PPTX-Erstellung fehlgeschlagen: {e}", exc_info=True)
            return None
    
    def _add_native_shapes_to_slide(self, slide, slide_spec: Dict[str, Any]):
        """
        Fügt native PowerPoint-Shapes zu einer Slide hinzu
        
        Diese Shapes sind editierbar in PowerPoint (nicht nur Bilder)
        """
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN
        
        # Shape-Mapping: VDL-Namen -> MSO_SHAPE
        shape_map = {
            'rectangle': MSO_SHAPE.RECTANGLE,
            'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
            'oval': MSO_SHAPE.OVAL,
            'circle': MSO_SHAPE.OVAL,
            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
            'diamond': MSO_SHAPE.DIAMOND,
            'pentagon': MSO_SHAPE.PENTAGON,
            'hexagon': MSO_SHAPE.HEXAGON,
            'octagon': MSO_SHAPE.OCTAGON,
            
            # Pfeile
            'right_arrow': MSO_SHAPE.RIGHT_ARROW,
            'left_arrow': MSO_SHAPE.LEFT_ARROW,
            'up_arrow': MSO_SHAPE.UP_ARROW,
            'down_arrow': MSO_SHAPE.DOWN_ARROW,
            'left_right_arrow': MSO_SHAPE.LEFT_RIGHT_ARROW,
            'circular_arrow': MSO_SHAPE.CIRCULAR_ARROW,
            'bent_arrow': MSO_SHAPE.BENT_ARROW,
            
            # Flowchart
            'flowchart_process': MSO_SHAPE.FLOWCHART_PROCESS,
            'flowchart_decision': MSO_SHAPE.FLOWCHART_DECISION,
            'flowchart_terminator': MSO_SHAPE.FLOWCHART_TERMINATOR,
            'flowchart_data': MSO_SHAPE.FLOWCHART_DATA,
            'flowchart_document': MSO_SHAPE.FLOWCHART_DOCUMENT,
            'flowchart_database': MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,  # Note: Using magnetic disk shape as database representation
            
            # Callouts
            'cloud_callout': MSO_SHAPE.CLOUD_CALLOUT,
            'oval_callout': MSO_SHAPE.OVAL_CALLOUT,
            
            # Sterne
            'star_5_point': MSO_SHAPE.STAR_5_POINT,
            'star_6_point': MSO_SHAPE.STAR_6_POINT,
        }
        
        for element in slide_spec.get('elements', []):
            element_type = element.get('type')
            
            if element_type == 'text':
                # Textbox hinzufügen
                pos = element.get('position', {'x': 0, 'y': 0})
                size = element.get('size', {'width': 400, 'height': 100})
                props = element.get('properties', {})
                
                textbox = slide.shapes.add_textbox(
                    Inches(pos['x'] / 100),
                    Inches(pos['y'] / 100),
                    Inches(size['width'] / 100),
                    Inches(size['height'] / 100)
                )
                textbox.text = element.get('content', '')
                
                # Formatierung
                if props.get('font_size'):
                    textbox.text_frame.paragraphs[0].font.size = Pt(props['font_size'])
                if props.get('color'):
                    color_hex = props['color'].lstrip('#')
                    r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                    textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
                if props.get('align'):
                    align_map = {
                        'left': PP_ALIGN.LEFT,
                        'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT
                    }
                    textbox.text_frame.paragraphs[0].alignment = align_map.get(props['align'], PP_ALIGN.LEFT)
            
            elif element_type == 'shape':
                # Shape hinzufügen
                shape_type = element.get('shape', 'rectangle')
                pos = element.get('position', {'x': 0, 'y': 0})
                size = element.get('size', {'width': 100, 'height': 100})
                props = element.get('properties', {})
                
                if shape_type in shape_map:
                    shape = slide.shapes.add_shape(
                        shape_map[shape_type],
                        Inches(pos['x'] / 100),
                        Inches(pos['y'] / 100),
                        Inches(size['width'] / 100),
                        Inches(size['height'] / 100)
                    )
                    
                    # Text in Shape
                    if 'content' in element and element['content']:
                        shape.text = element['content']
                    
                    # Füllfarbe
                    if props.get('fill_color'):
                        shape.fill.solid()
                        color_hex = props['fill_color'].lstrip('#')
                        r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                        shape.fill.fore_color.rgb = RGBColor(r, g, b)
                    
                    # Rahmen
                    if props.get('border_color'):
                        color_hex = props['border_color'].lstrip('#')
                        r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                        shape.line.color.rgb = RGBColor(r, g, b)
                    if props.get('border_width'):
                        shape.line.width = Pt(props['border_width'])
            
            elif element_type == 'connector':
                # Connector hinzufügen
                connector_type = element.get('connector_type', 'straight')
                start = element.get('start', {'x': 0, 'y': 0})
                end = element.get('end', {'x': 100, 'y': 100})
                
                connector_map = {
                    'straight': MSO_CONNECTOR.STRAIGHT,
                    'elbow': MSO_CONNECTOR.ELBOW,
                    'curve': MSO_CONNECTOR.CURVE
                }
                
                connector = slide.shapes.add_connector(
                    connector_map.get(connector_type, MSO_CONNECTOR.STRAIGHT),
                    Inches(start['x'] / 100),
                    Inches(start['y'] / 100),
                    Inches(end['x'] / 100),
                    Inches(end['y'] / 100)
                )
                
                # Connector-Formatierung
                props = element.get('properties', {})
                if props.get('width'):
                    connector.line.width = Pt(props['width'])
                if props.get('color'):
                    color_hex = props['color'].lstrip('#')
                    r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                    connector.line.color.rgb = RGBColor(r, g, b)



# Standalone-Test
if __name__ == '__main__':
    import asyncio
    
    async def test():
        agent = PresentationCanvasAgent()
        
        # Test: Präsentation generieren
        result = await agent.generate_presentation(
            "Erstelle eine Präsentation über BImSchG-Anlagen mit 2 Folien: "
            "Titel-Slide und eine Content-Slide mit Übersicht"
        )
        
        if result['success']:
            print(f"✅ Präsentation erstellt: {result['slide_count']} Folien")
            print(f"   VDL Slides: {len(result['vdl']['slides'])}")
            if result['pptx_path']:
                print(f"   PPTX: {result['pptx_path']}")
            for i, slide in enumerate(result['slides']):
                print(f"   Slide {i+1}: {slide['png_path']}")
        else:
            print(f"❌ Fehler: {result['error']}")
    
    asyncio.run(test())
